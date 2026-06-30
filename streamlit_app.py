import streamlit as st
import pandas as pd
import numpy as np
import sqlite3
from datetime import datetime
from pathlib import Path
import pickle
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
import warnings
import json

warnings.filterwarnings('ignore')

# ════════════════════════════════════════════════════════════════════════════
# CONFIGURATION
# ════════════════════════════════════════════════════════════════════════════

st.set_page_config(
    page_title="Candidate Dashboard Generator",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Colors
COLORS = {
    'red': '#FF000F',
    'lilac': '#6764f6',
    'red1': '#ff957e',
    'lil1': '#93a1ff',
    'red2': '#ffdccd',
    'dark': '#1a1a1a',
    'gray': '#888888',
}

CAND_COLORS = [COLORS['red'], COLORS['lilac'], COLORS['red1'], COLORS['lil1'], COLORS['red2']]

# ════════════════════════════════════════════════════════════════════════════
# DATABASE SETUP
# ════════════════════════════════════════════════════════════════════════════

def init_db():
    """Initialize SQLite database for storing dashboard history."""
    conn = sqlite3.connect('dashboards.db')
    c = conn.cursor()
    
    c.execute('''CREATE TABLE IF NOT EXISTS dashboards (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        name TEXT NOT NULL,
        description TEXT,
        created_at TIMESTAMP,
        updated_at TIMESTAMP,
        candidates_json TEXT,
        scores_json TEXT,
        data_json TEXT,
        version INTEGER,
        tags TEXT
    )''')
    
    conn.commit()
    return conn

@st.cache_resource
def get_db_connection():
    """Get database connection (cached)."""
    return init_db()

def save_dashboard_to_db(name, description, candidates, all_data, tags=""):
    """Save dashboard configuration to database."""
    conn = get_db_connection()
    c = conn.cursor()
    
    # Serialize data
    candidates_json = json.dumps(candidates)
    scores_json = json.dumps({c: all_data[c]['total_score'] for c in candidates})
    data_json = json.dumps({
        c: {
            'total': all_data[c]['total_score'],
            'interview': all_data[c]['interview_total'],
            'competency': all_data[c]['competency_total'],
        }
        for c in candidates
    })
    
    now = datetime.now().isoformat()
    
    c.execute('''INSERT INTO dashboards 
        (name, description, created_at, updated_at, candidates_json, scores_json, data_json, version, tags)
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
    ''', (name, description, now, now, candidates_json, scores_json, data_json, 1, tags))
    
    conn.commit()
    return c.lastrowid

def load_dashboards_from_db():
    """Load all saved dashboards from database."""
    conn = get_db_connection()
    c = conn.cursor()
    
    c.execute('''SELECT id, name, description, created_at, candidates_json, tags 
               FROM dashboards ORDER BY created_at DESC''')
    
    results = c.fetchall()
    return results

def load_dashboard_by_id(dashboard_id):
    """Load specific dashboard by ID."""
    conn = get_db_connection()
    c = conn.cursor()
    
    c.execute('''SELECT * FROM dashboards WHERE id = ?''', (dashboard_id,))
    result = c.fetchone()
    
    if result:
        return {
            'id': result[0],
            'name': result[1],
            'description': result[2],
            'created_at': result[3],
            'updated_at': result[4],
            'candidates_json': result[5],
            'scores_json': result[6],
            'data_json': result[7],
            'version': result[8],
            'tags': result[9]
        }
    return None

def delete_dashboard(dashboard_id):
    """Delete a saved dashboard."""
    conn = get_db_connection()
    c = conn.cursor()
    
    c.execute('''DELETE FROM dashboards WHERE id = ?''', (dashboard_id,))
    conn.commit()

# ════════════════════════════════════════════════════════════════════════════
# DATA EXTRACTION
# ════════════════════════════════════════════════════════════════════════════

STANDARD_CAT_ROWS = {
    13: 'Communications / Interpersonal',
    19: 'Functional / Technical Skills',
    25: 'Leadership / Initiative',
    31: 'Teamwork',
    37: 'Continual Learning',
    44: 'Diversity, Equity & Inclusion',
    50: 'Interpersonal / Change Mgmt',
}

STEPHEN_CAT_ROWS = {
    10: 'Communications / Interpersonal',
    16: 'Functional / Technical Skills',
    22: 'Leadership / Initiative',
    28: 'Teamwork',
    34: 'Interpersonal / Change Mgmt',
    40: 'Diversity, Equity & Inclusion',
    46: 'Continual Learning',
}

CAT_NAMES_STD = [
    'Communications / Interpersonal',
    'Functional / Technical Skills',
    'Leadership / Initiative',
    'Teamwork',
    'Continual Learning',
    'Diversity, Equity & Inclusion',
    'Interpersonal / Change Mgmt',
]

def extract_candidate_data(df, filename, is_stephen=False):
    """Extract scores from Excel file."""
    try:
        if is_stephen:
            total_score = df.iloc[1, 8]
            interview_total = df.iloc[8, 9]
            competency_total = df.iloc[8, 15]
            cat_rows = STEPHEN_CAT_ROWS
            score_col = 9
        else:
            total_score = df.iloc[2, 8]
            interview_total = df.iloc[10, 9]
            competency_total = df.iloc[10, 15]
            cat_rows = STANDARD_CAT_ROWS
            score_col = 9
        
        cat_scores = {}
        for row_idx, cat_name in cat_rows.items():
            score = df.iloc[row_idx, score_col]
            cat_scores[cat_name] = float(score) if pd.notna(score) else np.nan
        
        competencies = {}
        for idx in range(10, len(df)):
            row = df.iloc[idx]
            comp_name = row[10] if pd.notna(row[10]) else None
            comp_score = row[16] if pd.notna(row[16]) else None
            if comp_name and comp_score:
                competencies[comp_name.strip()] = float(comp_score)
        
        return {
            'total_score': float(total_score),
            'interview_total': float(interview_total),
            'competency_total': float(competency_total),
            'categories': cat_scores,
            'competencies': competencies,
        }
    except Exception as e:
        return None

def process_uploaded_files(uploaded_files):
    """Process uploaded Excel files."""
    all_data = {}
    
    for uploaded_file in uploaded_files:
        try:
            df = pd.read_excel(uploaded_file, sheet_name='Interview Tool', header=None)
            is_stephen = 'stephen' in uploaded_file.name.lower()
            
            # Extract candidate name from filename
            name_parts = uploaded_file.name.replace('.xlsx', '').split(' - ')
            if len(name_parts) >= 2:
                candidate_name = name_parts[1]
            else:
                candidate_name = uploaded_file.name
            
            data = extract_candidate_data(df, uploaded_file.name, is_stephen=is_stephen)
            if data:
                all_data[candidate_name] = data
        except Exception as e:
            st.warning(f"Error processing {uploaded_file.name}: {str(e)}")
    
    return all_data

# ════════════════════════════════════════════════════════════════════════════
# CHART GENERATION
# ════════════════════════════════════════════════════════════════════════════

def create_summary_table(candidates, total_scores):
    """Create summary DataFrame."""
    sorted_idx = np.argsort(total_scores)[::-1]
    sorted_cands = [candidates[i] for i in sorted_idx]
    sorted_scores = [total_scores[i] for i in sorted_idx]
    sorted_interview = [st.session_state.all_data[c]['interview_total'] for c in sorted_cands]
    sorted_competency = [st.session_state.all_data[c]['competency_total'] for c in sorted_cands]
    
    df = pd.DataFrame({
        'Rank': [f'#{i+1}' for i in range(len(sorted_cands))],
        'Candidate': sorted_cands,
        'Total Score': sorted_scores,
        'Interview Section': sorted_interview,
        'Competency Section': sorted_competency,
    })
    
    return df

def plot_total_scores(candidates, total_scores):
    """Plot total scores."""
    sorted_idx = np.argsort(total_scores)[::-1]
    sorted_cands = [candidates[i] for i in sorted_idx]
    sorted_scores = [total_scores[i] for i in sorted_idx]
    sorted_colors = [CAND_COLORS[i] for i in sorted_idx]
    
    fig, ax = plt.subplots(figsize=(10, len(candidates)*0.8))
    fig.patch.set_facecolor('white')
    
    bars = ax.barh(sorted_cands, sorted_scores, color=sorted_colors, height=0.5, zorder=3)
    
    for bar, score in zip(bars, sorted_scores):
        ax.text(score + 2, bar.get_y() + bar.get_height()/2, f'{score}',
                va='center', ha='left', fontsize=11, fontweight='bold')
    
    ax.set_xlim(0, max(total_scores) * 1.1)
    ax.set_xlabel('Total Score', fontsize=11)
    ax.set_title('Overall Candidate Rankings', fontsize=13, fontweight='bold', pad=15)
    ax.set_facecolor('#fafafa')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.grid(axis='x', linestyle='--', alpha=0.3, zorder=1)
    
    plt.tight_layout()
    return fig

def plot_stacked_scores(candidates, interview_totals, competency_totals, total_scores):
    """Plot stacked bar chart."""
    fig, ax = plt.subplots(figsize=(10, 5))
    fig.patch.set_facecolor('white')
    
    x = np.arange(len(candidates))
    bar_w = 0.5
    
    ax.bar(x, interview_totals, bar_w, label='Interview Section', color=COLORS['red'], zorder=3)
    ax.bar(x, competency_totals, bar_w, bottom=interview_totals, 
           label='Competency Section', color=COLORS['lilac'], zorder=3, alpha=0.85)
    
    for i, (iv, cv, tot) in enumerate(zip(interview_totals, competency_totals, total_scores)):
        ax.text(i, iv / 2, f'{iv}', ha='center', va='center', fontsize=9.5, 
                fontweight='bold', color='white')
        ax.text(i, iv + cv / 2, f'{cv}', ha='center', va='center', fontsize=9.5, 
                fontweight='bold', color='white')
        ax.text(i, tot + 4, f'{tot}', ha='center', va='bottom', fontsize=10, fontweight='bold')
    
    ax.set_xticks(x)
    ax.set_xticklabels(candidates, fontsize=10)
    ax.set_ylim(0, max(total_scores) * 1.2)
    ax.set_ylabel('Score', fontsize=11)
    ax.set_title('Score Composition: Interview + Competency', fontsize=12, fontweight='bold')
    ax.set_facecolor('#fafafa')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    ax.grid(axis='y', linestyle='--', alpha=0.3, zorder=1)
    ax.legend(loc='upper left', fontsize=9)
    
    plt.tight_layout()
    return fig

# ════════════════════════════════════════════════════════════════════════════
# PRESENTATION GENERATION
# ════════════════════════════════════════════════════════════════════════════

def generate_powerpoint(candidates, all_data):
    """Generate PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Colors for PowerPoint
    RGB_RED = RGBColor(0xFF, 0x00, 0x0F)
    RGB_LIL = RGBColor(0x67, 0x64, 0xF6)
    RGB_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    RGB_DARK = RGBColor(0x1A, 0x1A, 0x1A)
    RGB_LIGHTGRAY = RGBColor(0xF5, 0xF5, 0xF5)
    
    total_scores = [all_data[c]['total_score'] for c in candidates]
    interview_totals = [all_data[c]['interview_total'] for c in candidates]
    competency_totals = [all_data[c]['competency_total'] for c in candidates]
    
    # SLIDE 1 - Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    
    p = title_frame.paragraphs[0]
    p.text = "Candidate Interview Dashboard"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGB_RED
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = f"{len(candidates)} Candidates Evaluated | {datetime.now().strftime('%B %d, %Y')}"
    p.font.size = Pt(20)
    p.font.color.rgb = RGB_DARK
    
    # SLIDE 2 - Summary Table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Candidate Summary"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGB_DARK
    
    # Create table
    sorted_idx = np.argsort(total_scores)[::-1]
    sorted_cands = [candidates[i] for i in sorted_idx]
    sorted_scores = [total_scores[i] for i in sorted_idx]
    sorted_interview = [interview_totals[i] for i in sorted_idx]
    sorted_competency = [competency_totals[i] for i in sorted_idx]
    
    table_data = [['Rank', 'Candidate', 'Total Score', 'Interview', 'Competency']]
    for rank, (cand, score, intv, comp) in enumerate(
        zip(sorted_cands, sorted_scores, sorted_interview, sorted_competency), 1):
        table_data.append([f'#{rank}', cand, str(score), str(intv), str(comp)])
    
    rows, cols = len(table_data), len(table_data[0])
    tbl = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.1), 
                                 Inches(9), Inches(5.5)).table
    
    for r_idx, row_data in enumerate(table_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = cell_text
            tf = cell.text_frame
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.CENTER
                for run in para.runs:
                    run.font.size = Pt(11)
                    run.font.bold = (r_idx == 0)
                    run.font.color.rgb = RGB_WHITE if r_idx == 0 else RGB_DARK
            
            if r_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGB_RED
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGB_LIGHTGRAY if r_idx % 2 == 0 else RGB_WHITE
    
    return prs

# ════════════════════════════════════════════════════════════════════════════
# MAIN APP
# ════════════════════════════════════════════════════════════════════════════

# Initialize session state
if 'all_data' not in st.session_state:
    st.session_state.all_data = {}
if 'candidates' not in st.session_state:
    st.session_state.candidates = []

# ════════════════════════════════════════════════════════════════════════════

st.title("📊 Candidate Dashboard Generator")
st.markdown("Professional hiring dashboards with history & versioning")

# ════════════════════════════════════════════════════════════════════════════

tab1, tab2, tab3 = st.tabs(["🆕 New Dashboard", "📚 History", "📈 Compare"])

# ════════════════════════════════════════════════════════════════════════════
# TAB 1: NEW DASHBOARD
# ════════════════════════════════════════════════════════════════════════════

with tab1:
    st.header("Create New Dashboard")
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.subheader("📤 Upload Evaluation Files")
        uploaded_files = st.file_uploader(
            "Drag and drop or click to select Excel files",
            type="xlsx",
            accept_multiple_files=True,
            key="file_uploader"
        )
    
    with col2:
        st.subheader("ℹ️ File Format")
        st.info("""
        **Expected filename:**
        ```
        Candidate Evaluation Form 
        - [Name] - [Score].xlsx
        ```
        
        **Example:**
        ```
        Candidate Evaluation Form 
        - John Smith - 150.xlsx
        ```
        """)
    
    if uploaded_files:
        st.success(f"✅ {len(uploaded_files)} file(s) selected")
        
        # Process files
        st.session_state.all_data = process_uploaded_files(uploaded_files)
        st.session_state.candidates = sorted(st.session_state.all_data.keys())
        
        if st.session_state.candidates:
            st.subheader("👥 Preview Candidates")
            
            # Show summary
            total_scores = [st.session_state.all_data[c]['total_score'] for c in st.session_state.candidates]
            summary_df = create_summary_table(st.session_state.candidates, total_scores)
            
            st.dataframe(summary_df, use_container_width=True, hide_index=True)
            
            # Show charts
            col1, col2 = st.columns(2)
            
            with col1:
                st.pyplot(plot_total_scores(st.session_state.candidates, total_scores))
            
            with col2:
                interview_totals = [st.session_state.all_data[c]['interview_total'] 
                                   for c in st.session_state.candidates]
                competency_totals = [st.session_state.all_data[c]['competency_total'] 
                                    for c in st.session_state.candidates]
                st.pyplot(plot_stacked_scores(st.session_state.candidates, 
                                             interview_totals, competency_totals, total_scores))
            
            # Save and generate options
            st.divider()
            st.subheader("💾 Save & Generate")
            
            col1, col2, col3 = st.columns(3)
            
            with col1:
                dashboard_name = st.text_input("Dashboard Name", 
                    f"Round_{datetime.now().strftime('%Y%m%d_%H%M')}")
            
            with col2:
                dashboard_tags = st.text_input("Tags (comma-separated)", "hiring,2026")
            
            with col3:
                dashboard_description = st.text_area("Description (optional)", height=80)
            
            col1, col2 = st.columns(2)
            
            with col1:
                if st.button("💾 Save Dashboard to History", use_container_width=True):
                    dashboard_id = save_dashboard_to_db(
                        dashboard_name, 
                        dashboard_description, 
                        st.session_state.candidates, 
                        st.session_state.all_data,
                        dashboard_tags
                    )
                    st.success(f"✅ Dashboard saved! ID: {dashboard_id}")
            
            with col2:
                if st.button("📥 Generate & Download PowerPoint", use_container_width=True, type="primary"):
                    with st.spinner("Generating dashboard..."):
                        prs = generate_powerpoint(st.session_state.candidates, st.session_state.all_data)
                        
                        # Convert to bytes
                        pptx_bytes = io.BytesIO()
                        prs.save(pptx_bytes)
                        pptx_bytes.seek(0)
                        
                        # Also save to history
                        dashboard_id = save_dashboard_to_db(
                            dashboard_name, 
                            dashboard_description, 
                            st.session_state.candidates, 
                            st.session_state.all_data,
                            dashboard_tags
                        )
                        
                        st.success("✅ Dashboard generated!")
                        
                        st.download_button(
                            label="📥 Download PowerPoint",
                            data=pptx_bytes,
                            file_name=f"{dashboard_name}.pptx",
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            use_container_width=True
                        )

# ════════════════════════════════════════════════════════════════════════════
# TAB 2: HISTORY
# ════════════════════════════════════════════════════════════════════════════

with tab2:
    st.header("📚 Dashboard History")
    
    dashboards = load_dashboards_from_db()
    
    if dashboards:
        st.info(f"📊 {len(dashboards)} dashboard(s) saved")
        
        for db_id, name, description, created_at, candidates_json, tags in dashboards:
            with st.expander(f"📋 {name} — {created_at[:10]}", expanded=False):
                col1, col2, col3 = st.columns([2, 1, 1])
                
                with col1:
                    st.write(f"**Description:** {description or 'None'}")
                    st.write(f"**Tags:** {tags or 'None'}")
                    st.write(f"**Created:** {created_at}")
                    candidates_list = json.loads(candidates_json)
                    st.write(f"**Candidates:** {len(candidates_list)} ({', '.join(candidates_list[:3])}...)")
                
                with col2:
                    if st.button("🔄 Reload", key=f"reload_{db_id}"):
                        dashboard = load_dashboard_by_id(db_id)
                        st.session_state.candidates = json.loads(dashboard['candidates_json'])
                        
                        # Reconstruct all_data from scores
                        scores = json.loads(dashboard['scores_json'])
                        data = json.loads(dashboard['data_json'])
                        
                        st.session_state.all_data = {
                            c: {
                                'total_score': scores[c],
                                'interview_total': data[c]['interview'],
                                'competency_total': data[c]['competency'],
                                'categories': {},
                                'competencies': {}
                            }
                            for c in st.session_state.candidates
                        }
                        
                        st.success("✅ Loaded! Go to 'New Dashboard' tab")
                        st.rerun()
                
                with col3:
                    if st.button("🗑️ Delete", key=f"delete_{db_id}"):
                        delete_dashboard(db_id)
                        st.success("✅ Deleted")
                        st.rerun()
    else:
        st.info("📭 No saved dashboards yet. Create one in the 'New Dashboard' tab!")

# ════════════════════════════════════════════════════════════════════════════
# TAB 3: COMPARE
# ════════════════════════════════════════════════════════════════════════════

with tab3:
    st.header("📈 Compare Dashboards")
    
    dashboards = load_dashboards_from_db()
    
    if len(dashboards) >= 2:
        dashboard_options = {db[1]: db for db in dashboards}
        
        col1, col2 = st.columns(2)
        
        with col1:
            dash1_name = st.selectbox("Dashboard 1", list(dashboard_options.keys()))
        
        with col2:
            dash2_name = st.selectbox("Dashboard 2", list(dashboard_options.keys()), key="dash2")
        
        if dash1_name and dash2_name and dash1_name != dash2_name:
            dash1 = dashboard_options[dash1_name]
            dash2 = dashboard_options[dash2_name]
            
            candidates1 = json.loads(dash1[4])
            candidates2 = json.loads(dash2[4])
            scores1 = json.loads(dash1[5])
            scores2 = json.loads(dash2[5])
            
            st.subheader("Comparison Summary")
            
            comparison_data = []
            all_candidates = set(candidates1 + candidates2)
            
            for cand in sorted(all_candidates):
                score1 = scores1.get(cand, "-")
                score2 = scores2.get(cand, "-")
                diff = score2 - score1 if isinstance(score1, (int, float)) and isinstance(score2, (int, float)) else "-"
                
                comparison_data.append({
                    'Candidate': cand,
                    'Score (Dash 1)': score1,
                    'Score (Dash 2)': score2,
                    'Change': diff
                })
            
            comparison_df = pd.DataFrame(comparison_data)
            st.dataframe(comparison_df, use_container_width=True, hide_index=True)
            
            st.info("💡 Positive change = improvement, Negative = decline")
    else:
        st.info("📭 Need at least 2 saved dashboards to compare. Create more dashboards first!")

# ════════════════════════════════════════════════════════════════════════════

st.divider()
st.markdown("""
---
**Candidate Dashboard Generator** | Cloud-based hiring analytics  
Powered by Streamlit | Data secured locally  
Questions? Contact: chad.brant@us.abb.com
""")
